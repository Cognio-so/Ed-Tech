import os
import time
import json
import logging
import mimetypes
from typing import List, Dict, Optional, Tuple
import io
import asyncio

import requests
from dotenv import load_dotenv
from pptx import Presentation
from openai import OpenAI

try:
    import cloudinary
    import cloudinary.uploader
    from cloudinary.utils import cloudinary_url
except ImportError:
    raise ImportError("Cloudinary SDK not found. Install with: pip install cloudinary")


try:
    import asposeslidescloud
    from asposeslidescloud.api_client import ApiClient
    from asposeslidescloud.configuration import Configuration
    from asposeslidescloud.apis.slides_api import SlidesApi
except ImportError:
    raise ImportError("Aspose.Slides Cloud SDK not found. Install with: pip install asposeslidescloud")


load_dotenv()

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


class CloudinaryStorage:
    """
    A storage manager to interact with Cloudinary for retrieving and uploading files.
    """
    def __init__(self):
        cloud_name = os.getenv("CLOUDINARY_CLOUD_NAME")
        api_key = os.getenv("CLOUDINARY_API_KEY")
        api_secret = os.getenv("CLOUDINARY_API_SECRET")

        if not all([cloud_name, api_key, api_secret]):
            raise ValueError("Missing Cloudinary environment variables: CLOUDINARY_CLOUD_NAME, CLOUDINARY_API_KEY, CLOUDINARY_API_SECRET")

        cloudinary.config(
            cloud_name=cloud_name,
            api_key=api_key,
            api_secret=api_secret,
            secure=True
        )
        logging.info("Cloudinary storage manager initialized.")

    def get_file_content_bytes(self, public_id: str) -> Optional[bytes]:
        """
        Retrieves a file from Cloudinary using its public ID.
        For non-image files like PPTX, we must specify resource_type='raw'.
        """
        try:
            file_url, _ = cloudinary_url(public_id, resource_type="raw")
            logging.info(f"Fetching file from Cloudinary URL: {file_url}")
            response = requests.get(file_url, timeout=60)
            response.raise_for_status()
            return response.content
        except Exception as e:
            logging.error(f"Failed to retrieve file '{public_id}' from Cloudinary: {e}")
            return None

    def upload_file(self, file_data: bytes, public_id: str) -> Tuple[bool, str]:
        """
        Uploads a file to Cloudinary.
        Returns a tuple (success, public_id_or_error_message).
        """
        try:
            logging.info(f"Uploading file to Cloudinary with public_id: {public_id}")
            upload_result = cloudinary.uploader.upload(
                io.BytesIO(file_data),
                public_id=public_id,
                resource_type="raw",
                overwrite=True
            )

            returned_public_id = upload_result.get("public_id")
            if not returned_public_id:
                raise Exception("Cloudinary upload did not return a public_id.")
            return True, returned_public_id
        except Exception as e:
            error_msg = f"Cloudinary upload failed: {e}"
            logging.error(error_msg)
            return False, error_msg

    async def upload_file_async(self, file_data: bytes, public_id: str) -> Tuple[bool, str]:
        """Asynchronous wrapper for upload_file."""
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(
            None,
            self.upload_file,
            file_data,
            public_id
        )


class PPTXToHeyGenVideo:
    def __init__(
        self,
        storage_manager: CloudinaryStorage,
        heygen_api_key: Optional[str] = None,
        openai_api_key: Optional[str] = None,
        pptx_avatar_id: Optional[str] = None,
        pptx_voice_id: Optional[str] = None,
        language: str = "english",
        model: str = "gpt-4o",
        width: int = 1280,
        height: int = 720,
        request_timeout_s: int = 90,
        poll_interval_s: int = 10,
    ):

        self.storage_manager = storage_manager
        self.heygen_api_key = heygen_api_key or os.getenv("HEYGEN_API_KEY")
        self.openai_api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
        self.aspose_client_id = os.getenv("ASPOSE_CLIENT_ID")
        self.aspose_client_secret = os.getenv("ASPOSE_CLIENT_SECRET")

        self.avatar_id = pptx_avatar_id or os.getenv("HEYGEN_AVATAR_ID")
        self.voice_id = pptx_voice_id or os.getenv("HEYGEN_VOICE_ID")
        self.language = language
        self.model = model
        self.width = width
        self.height = height
        self.background_color = "#FFFFFF"
        self.request_timeout_s = request_timeout_s
        self.poll_interval_s = poll_interval_s
        self.slide_asset_ids: List[str] = []

        if not self.storage_manager:
            raise ValueError("A storage_manager instance is required.")
        if not self.heygen_api_key:
            raise ValueError("Missing HEYGEN_API_KEY")
        if not self.openai_api_key:
            raise ValueError("Missing OPENAI_API_KEY")
        if not self.aspose_client_id or not self.aspose_client_secret:
            raise ValueError("Missing ASPOSE_CLIENT_ID or ASPOSE_CLIENT_SECRET")
        if not self.avatar_id:
            raise ValueError("Missing HEYGEN_AVATAR_ID")
        if not self.voice_id:
            raise ValueError("Missing HEYGEN_VOICE_ID")

        self._openai_client = OpenAI(api_key=self.openai_api_key)
        self._heygen_base_v2 = "https://api.heygen.com/v2"
        self._headers = {
            "Accept": "application/json",
            "Content-Type": "application/json",
            "X-Api-Key": self.heygen_api_key,
        }
        
        aspose_config = Configuration()
        aspose_config.app_sid = self.aspose_client_id
        aspose_config.app_key = self.aspose_client_secret
        self._slides_api = SlidesApi(aspose_config)

    def _post_with_retry(self, url: str, payload: Dict) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.post(url, headers=self._headers, json=payload, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"POST attempt {attempt + 1} failed: {e}")
                if attempt == 2:
                    raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    def _get_with_retry(self, url: str) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.get(url, headers=self._headers, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"GET attempt {attempt + 1} failed: {e}")
                if attempt == 2:
                    raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    def _upload_asset_to_heygen(self, file_bytes: bytes, file_name: str) -> str:
        mime_type, _ = mimetypes.guess_type(file_name)
        if not mime_type:
            mime_type = "image/png"

        logging.info(f"Uploading {file_name} to HeyGen...")
        url = "https://upload.heygen.com/v1/asset"
        upload_headers = {"X-Api-Key": self.heygen_api_key, "Content-Type": mime_type}
        file_like_object = io.BytesIO(file_bytes)

        for attempt in range(3):
            try:
                file_like_object.seek(0)
                resp = requests.post(url, headers=upload_headers, data=file_like_object, timeout=self.request_timeout_s)
                resp.raise_for_status()
                response_data = resp.json()
                asset_id = response_data.get("data", {}).get("id")

                if not asset_id:
                    raise RuntimeError(f"Failed to get asset_id from HeyGen response: {response_data}")

                logging.info(f"Successfully uploaded {file_name}. Asset ID: {asset_id}")
                return asset_id
            except requests.RequestException as e:
                logging.warning(f"Upload attempt {attempt + 1} for {file_name} failed: {e}")
                if attempt == 2:
                    raise
                time.sleep(2 ** attempt)
        raise RuntimeError(f"File upload failed for {file_name} after multiple retries.")

    def _convert_pptx_and_upload_slides(self, pptx_storage_key: str, pptx_bytes: bytes, num_slides: int) -> List[str]:
        temp_storage_keys = []
        logging.info(f"Converting {num_slides} slides to images using Aspose.Slides Cloud...")
        for i in range(1, num_slides + 1):
            temp_file_path = None
            try:
                temp_file_path = self._slides_api.download_slide_online(
                    document=io.BytesIO(pptx_bytes),
                    slide_index=i,
                    format='PNG'
                )
                
                # We must now open and read the bytes from this temporary file.
                with open(temp_file_path, 'rb') as f:
                    slide_image_bytes = f.read()

                if not slide_image_bytes:
                    raise ValueError(f"Aspose API returned an empty image for slide {i}.")
                
                slide_filename = f"slide_{i}.png"
                asset_id = self._upload_asset_to_heygen(slide_image_bytes, slide_filename)
                self.slide_asset_ids.append(asset_id)
                logging.info(f"Processed and uploaded slide {i}/{num_slides}.")

            except Exception as e:
                logging.error(f"Failed to process slide {i}: {e}", exc_info=True)
                continue
            finally:
                if temp_file_path and os.path.exists(temp_file_path):
                    try:
                        os.remove(temp_file_path)
                        logging.info(f"Cleaned up temporary file: {temp_file_path}")
                    except OSError as e:
                        logging.error(f"Error removing temporary file {temp_file_path}: {e}")

        return temp_storage_keys


    def _build_video_inputs(self, slide_notes: List[str]) -> List[Dict]:
        video_inputs: List[Dict] = []
        for idx, note in enumerate(slide_notes):
            scene = {
                "character": {
                    "type": "talking_photo",
                    "talking_photo_id": self.avatar_id,
                    "talking_photo_style": "circle",
                    "scale": 0.33,
                    "offset": {"x": 0.42, "y": 0.42},
                },
                "voice": {"type": "text", "input_text": note, "voice_id": self.voice_id},
            }
            if idx < len(self.slide_asset_ids):
                asset_id = self.slide_asset_ids[idx]
                scene["background"] = {"type": "image", "image_asset_id": asset_id}
            else:
                scene["background"] = {"type": "color", "value": self.background_color}
            video_inputs.append(scene)
        return video_inputs

    def create_multi_scene_video(self, slide_notes: List[str], title: str) -> Dict:
        video_inputs = self._build_video_inputs(slide_notes)
        payload = {
            "video_inputs": video_inputs,
            "dimension": {"width": self.width, "height": self.height},
            "title": title,
        }
        generate_url = f"{self._heygen_base_v2}/video/generate"
        resp = self._post_with_retry(generate_url, payload)
        data = resp.json()
        video_id = data.get("data", {}).get("video_id")
        logging.info(f"Generating video with video_id: {video_id}, Avatar ID: {self.avatar_id}, Voice ID: {self.voice_id}.")
        if not video_id:
            raise RuntimeError(f"Unexpected response from generate endpoint: {resp.text}")
        return {"video_id": video_id}

    def wait_for_video(self, video_id: str) -> Dict:
        logging.info(f"Waiting for video {video_id} to complete. Polling status...")
        while True:
            status_url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"
            try:
                resp = self._get_with_retry(status_url)
                data = resp.json().get("data", {})
                status = data.get("status")
                if status in ("completed", "success"):
                    logging.info("Video generation completed successfully.")
                    return {"status": "completed", "video_url": data.get("video_url")}
                elif status in ("failed", "error"):
                    logging.error(f"Video generation failed with data: {data}")
                    raise RuntimeError(f"Video generation failed: {data}")
                else:
                    logging.info(f"Video status is '{status}', continuing to wait...")
            except requests.RequestException as e:
                logging.warning(f"Status check failed due to a network error: {e}. Retrying...")
            time.sleep(self.poll_interval_s)
    
    def generate_speaker_notes(self, slide_texts: List[str]) -> List[str]:
        logging.info(f"Generating speaker notes for {len(slide_texts)} slides using model {self.model}...")
        speaker_notes = []
        for i, text in enumerate(slide_texts):
            try:
                system_prompt = (
                    f"You are a virtual teacher. Your task is to explain the content of a presentation slide in {self.language}. "
                    "You will be given the text from a slide and must generate a clear, engaging, and educational script. "
                    "Your tone should be professional yet approachable. Your entire response must be in {self.language}."
                )
                user_prompt = f"Here is the slide content:\n\n---\n\n{text}\n\n---\n\nPlease provide the speaker notes in {self.language}."
                
                response = self._openai_client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    temperature=0.7, max_tokens=300,
                )
                note = response.choices[0].message.content.strip()
                speaker_notes.append(note)
            except Exception as e:
                logging.error(f"Failed to generate notes for slide {i+1}: {e}")
                speaker_notes.append(f"Content for slide {i+1}: {text}")
        return speaker_notes

    def convert(self, pptx_public_id: str, title: Optional[str] = None, max_slides: Optional[int] = None) -> Dict:
        logging.info(f"Starting conversion for PPTX in Cloudinary with public_id: {pptx_public_id}")
        try:
            pptx_bytes = self.storage_manager.get_file_content_bytes(pptx_public_id)
            if not pptx_bytes:
                raise FileNotFoundError(f"Could not retrieve PPTX file from Cloudinary with public_id: {pptx_public_id}")
            prs = Presentation(io.BytesIO(pptx_bytes))
            slides = list(prs.slides)
            if not slides:
                raise ValueError("No slides found in the PowerPoint file.")
            if max_slides and len(slides) > max_slides:
                slides = slides[:max_slides]
            slides_text = ["\n".join(s.text for s in slide.shapes if hasattr(s, "text") and s.text).strip() for slide in slides]
            self._convert_pptx_and_upload_slides(pptx_public_id, pptx_bytes, len(slides))
            if len(self.slide_asset_ids) < len(slides_text):
                slides_text = slides_text[:len(self.slide_asset_ids)]
            slide_notes = self.generate_speaker_notes(slides_text)
            final_title = title or os.path.basename(pptx_public_id)
            gen_info = self.create_multi_scene_video(slide_notes, title=final_title)
            final_status = self.wait_for_video(gen_info["video_id"])
            return {
                "video_id": gen_info["video_id"],
                "video_url": final_status.get("video_url"),
                "slides_count": len(slide_notes),
            }
        finally:
            pass

if __name__ == "__main__":
    try:
        storage = CloudinaryStorage()
        pptx_public_id_in_cloudinary = "test"
        video_title = "My Test Presentation Video"
        if not storage.get_file_content_bytes(pptx_public_id_in_cloudinary):
             logging.error(f"File not found in Cloudinary storage with public_id: '{pptx_public_id_in_cloudinary}'")
        else:
            converter = PPTXToHeyGenVideo(storage_manager=storage)
            result = converter.convert(pptx_public_id_in_cloudinary, title=video_title)
            print("\n--- Conversion Successful ---")
            print(json.dumps(result, indent=2))
    except Exception as e:
        logging.error(f"An error occurred during conversion: {e}", exc_info=True)