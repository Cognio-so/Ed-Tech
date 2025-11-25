import os
import time
import json
import logging
import mimetypes
from typing import List, Dict, Optional, Tuple
import io
import asyncio

import requests
from dotenv import load_dotenv
from pptx import Presentation
from openai import OpenAI

try:
    import cloudinary
    import cloudinary.uploader
    import cloudinary.api
    from cloudinary.utils import cloudinary_url
except ImportError:
    raise ImportError("Cloudinary SDK not found. Install with: pip install cloudinary")

try:
    import asposeslidescloud
    from asposeslidescloud.api_client import ApiClient
    from asposeslidescloud.configuration import Configuration
    from asposeslidescloud.apis.slides_api import SlidesApi
except ImportError:
    raise ImportError("Aspose.Slides Cloud SDK not found. Install with: pip install asposeslidescloud")


load_dotenv()

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


class CloudinaryStorage:
    """
    A production-ready storage manager for Cloudinary.
    """
    def __init__(self):
        cloud_name = os.getenv("CLOUDINARY_CLOUD_NAME")
        api_key = os.getenv("CLOUDINARY_API_KEY")
        api_secret = os.getenv("CLOUDINARY_API_SECRET")

        if not all([cloud_name, api_key, api_secret]):
            raise ValueError("Missing Cloudinary environment variables.")

        cloudinary.config(
            cloud_name=cloud_name,
            api_key=api_key,
            api_secret=api_secret,
            secure=True  # Force HTTPS
        )
        logging.info("Cloudinary storage manager initialized.")

    def get_file_content_bytes(self, public_id: str) -> Optional[bytes]:
        """
        Retrieves a file from Cloudinary.
        Handles 'raw' resource types specifically for PPTX files.
        """
        # 1. Attempt to build the URL specifically for a RAW file (required for PPTX)
        # Note: For raw files, the public_id usually MUST include the extension (e.g. 'deck.pptx')
        file_url, options = cloudinary_url(public_id, resource_type="raw", secure=True)
        
        logging.info(f"Attempting to fetch file from Cloudinary: {file_url}")

        try:
            response = requests.get(file_url, timeout=60)
            
            # If 404, try checking if the ID needs an extension appended (common mistake)
            if response.status_code == 404 and not public_id.lower().endswith(".pptx"):
                logging.warning(f"404 encountered. Retrying with .pptx extension for public_id: {public_id}")
                return self.get_file_content_bytes(f"{public_id}.pptx")
                
            response.raise_for_status()
            return response.content
            
        except requests.RequestException as e:
            logging.error(f"Failed to retrieve file '{public_id}': {e}")
            return None

    def upload_file(self, file_path: str, public_id: str) -> Tuple[bool, str]:
        """
        Uploads a local file to Cloudinary as a RAW resource.
        """
        try:
            logging.info(f"Uploading local file '{file_path}' to Cloudinary as '{public_id}'...")
            
            upload_result = cloudinary.uploader.upload(
                file_path,
                public_id=public_id,
                resource_type="raw",
                overwrite=True,
                invalidate=True 
            )

            returned_public_id = upload_result.get("public_id")
            if not returned_public_id:
                raise Exception("Cloudinary upload did not return a public_id.")
            
            logging.info(f"Upload successful. Cloudinary Public ID: {returned_public_id}")
            return True, returned_public_id
        except Exception as e:
            error_msg = f"Cloudinary upload failed: {e}"
            logging.error(error_msg)
            return False, error_msg


class PPTXToHeyGenVideo:
    def __init__(
        self,
        storage_manager: CloudinaryStorage,
        heygen_api_key: Optional[str] = None,
        openai_api_key: Optional[str] = None,
        pptx_avatar_id: Optional[str] = None,
        pptx_voice_id: Optional[str] = None,
        language: str = "english",
        model: str = "gpt-4o",
        width: int = 1280,
        height: int = 720,
        request_timeout_s: int = 90,
        poll_interval_s: int = 10,
    ):

        self.storage_manager = storage_manager
        self.heygen_api_key = heygen_api_key or os.getenv("HEYGEN_API_KEY")
        self.openai_api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
        self.aspose_client_id = os.getenv("ASPOSE_CLIENT_ID")
        self.aspose_client_secret = os.getenv("ASPOSE_CLIENT_SECRET")

        self.avatar_id = pptx_avatar_id or os.getenv("HEYGEN_AVATAR_ID")
        self.voice_id = pptx_voice_id or os.getenv("HEYGEN_VOICE_ID")
        self.language = language
        self.model = model
        self.width = width
        self.height = height
        self.background_color = "#FFFFFF"
        self.request_timeout_s = request_timeout_s
        self.poll_interval_s = poll_interval_s
        self.slide_asset_ids: List[str] = []

        if not self.storage_manager: raise ValueError("Storage manager required.")
        if not self.heygen_api_key: raise ValueError("Missing HEYGEN_API_KEY")
        if not self.openai_api_key: raise ValueError("Missing OPENAI_API_KEY")
        if not self.aspose_client_id: raise ValueError("Missing ASPOSE_CLIENT_ID")
        if not self.aspose_client_secret: raise ValueError("Missing ASPOSE_CLIENT_SECRET")
        if not self.avatar_id: raise ValueError("Missing HEYGEN_AVATAR_ID")
        if not self.voice_id: raise ValueError("Missing HEYGEN_VOICE_ID")

        self._openai_client = OpenAI(api_key=self.openai_api_key)
        self._heygen_base_v2 = "https://api.heygen.com/v2"
        self._headers = {
            "Accept": "application/json",
            "Content-Type": "application/json",
            "X-Api-Key": self.heygen_api_key,
        }
        
        aspose_config = Configuration()
        aspose_config.app_sid = self.aspose_client_id
        aspose_config.app_key = self.aspose_client_secret
        self._slides_api = SlidesApi(aspose_config)

    def _post_with_retry(self, url: str, payload: Dict) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.post(url, headers=self._headers, json=payload, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"POST attempt {attempt + 1} failed: {e}")
                if attempt == 2: raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    def _get_with_retry(self, url: str) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.get(url, headers=self._headers, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"GET attempt {attempt + 1} failed: {e}")
                if attempt == 2: raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    def _upload_asset_to_heygen(self, file_bytes: bytes, file_name: str) -> str:
        mime_type, _ = mimetypes.guess_type(file_name)
        if not mime_type: mime_type = "image/png"

        logging.info(f"Uploading {file_name} to HeyGen...")
        url = "https://upload.heygen.com/v1/asset"
        upload_headers = {"X-Api-Key": self.heygen_api_key, "Content-Type": mime_type}
        file_like_object = io.BytesIO(file_bytes)

        for attempt in range(3):
            try:
                file_like_object.seek(0)
                resp = requests.post(url, headers=upload_headers, data=file_like_object, timeout=self.request_timeout_s)
                resp.raise_for_status()
                asset_id = resp.json().get("data", {}).get("id")
                if not asset_id: raise RuntimeError("No asset_id returned")
                return asset_id
            except requests.RequestException as e:
                logging.warning(f"HeyGen Upload attempt {attempt + 1} failed: {e}")
                if attempt == 2: raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Upload failed.")

    def _convert_pptx_and_upload_slides(self, pptx_storage_key: str, pptx_bytes: bytes, num_slides: int) -> List[str]:
        logging.info(f"Converting {num_slides} slides to images via Aspose...")
        # Aspose.Slides Cloud: download_slide_online returns a temporary local path
        for i in range(1, num_slides + 1):
            temp_file_path = None
            try:
                temp_file_path = self._slides_api.download_slide_online(
                    document=io.BytesIO(pptx_bytes),
                    slide_index=i,
                    format='PNG'
                )
                
                if not os.path.exists(temp_file_path):
                    raise FileNotFoundError(f"Aspose returned path {temp_file_path} but file is missing.")

                with open(temp_file_path, 'rb') as f:
                    slide_image_bytes = f.read()

                asset_id = self._upload_asset_to_heygen(slide_image_bytes, f"slide_{i}.png")
                self.slide_asset_ids.append(asset_id)
                logging.info(f"Uploaded slide {i}/{num_slides} (Asset: {asset_id}).")

            except Exception as e:
                logging.error(f"Failed to process slide {i}: {e}", exc_info=True)
            finally:
                if temp_file_path and os.path.exists(temp_file_path):
                    try:
                        os.remove(temp_file_path)
                    except OSError:
                        pass

        return self.slide_asset_ids

    def _build_video_inputs(self, slide_notes: List[str]) -> List[Dict]:
        video_inputs = []
        for idx, note in enumerate(slide_notes):
            scene = {
                "character": {
                    "type": "talking_photo",
                    "talking_photo_id": self.avatar_id,
                    "talking_photo_style": "circle",
                    "scale": 0.33,
                    "offset": {"x": 0.42, "y": 0.42},
                },
                "voice": {"type": "text", "input_text": note, "voice_id": self.voice_id},
            }
            if idx < len(self.slide_asset_ids):
                scene["background"] = {"type": "image", "image_asset_id": self.slide_asset_ids[idx]}
            else:
                scene["background"] = {"type": "color", "value": self.background_color}
            video_inputs.append(scene)
        return video_inputs

    def create_multi_scene_video(self, slide_notes: List[str], title: str) -> Dict:
        video_inputs = self._build_video_inputs(slide_notes)
        payload = {
            "video_inputs": video_inputs,
            "dimension": {"width": self.width, "height": self.height},
            "title": title,
        }
        generate_url = f"{self._heygen_base_v2}/video/generate"
        resp = self._post_with_retry(generate_url, payload)
        video_id = resp.json().get("data", {}).get("video_id")
        if not video_id:
            raise RuntimeError(f"No video_id in response: {resp.text}")
        return {"video_id": video_id}

    def wait_for_video(self, video_id: str) -> Dict:
        logging.info(f"Polling status for video {video_id}...")
        while True:
            status_url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"
            try:
                resp = self._get_with_retry(status_url)
                data = resp.json().get("data", {})
                status = data.get("status")
                
                if status == "completed":
                    logging.info("Video generation completed.")
                    return {"status": "completed", "video_url": data.get("video_url")}
                elif status == "failed":
                    raise RuntimeError(f"Video generation failed: {data.get('error')}")
                
                logging.info(f"Status: {status}. Waiting {self.poll_interval_s}s...")
            except requests.RequestException as e:
                logging.warning(f"Status check network error: {e}")
            
            time.sleep(self.poll_interval_s)

    def generate_speaker_notes(self, slide_texts: List[str]) -> List[str]:
        logging.info(f"Generating AI speaker notes for {len(slide_texts)} slides...")
        notes = []
        for i, text in enumerate(slide_texts):
            try:
                # Fallback if slide is empty
                if not text.strip():
                    text = "(Visual slide with no text)"

                system_prompt = (
                    f"You are a professional presenter. Summarize the following slide content into a spoken script "
                    f"in {self.language}. Keep it concise (under 30 seconds of speech)."
                )
                response = self._openai_client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": text},
                    ],
                    temperature=0.7, max_tokens=300,
                )
                notes.append(response.choices[0].message.content.strip())
            except Exception as e:
                logging.error(f"OpenAI error for slide {i+1}: {e}")
                notes.append(f"Slide {i+1}: {text[:100]}...")
        return notes

    def convert(self, pptx_public_id: str, title: Optional[str] = None, max_slides: Optional[int] = None) -> Dict:
        logging.info(f"Starting process for Cloudinary ID: {pptx_public_id}")
        
        # 1. Get PPTX
        pptx_bytes = self.storage_manager.get_file_content_bytes(pptx_public_id)
        if not pptx_bytes:
            raise FileNotFoundError(f"Could not download PPTX for ID: {pptx_public_id}. Check Cloudinary.")

        # 2. Parse PPTX
        prs = Presentation(io.BytesIO(pptx_bytes))
        slides = list(prs.slides)
        if max_slides: slides = slides[:max_slides]
        
        # Extract Text
        slides_text = []
        for slide in slides:
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)
            slides_text.append(" ".join(text_runs))

        # 3. Convert Slides to Images
        self._convert_pptx_and_upload_slides(pptx_public_id, pptx_bytes, len(slides))
        
        # 4. Generate Audio Script
        slide_notes = self.generate_speaker_notes(slides_text)

        # 5. Create Video
        final_title = title or f"Video_{pptx_public_id}"
        gen_info = self.create_multi_scene_video(slide_notes, title=final_title)
        
        # 6. Wait for result
        final_status = self.wait_for_video(gen_info["video_id"])
        
        return {
            "video_id": gen_info["video_id"],
            "video_url": final_status.get("video_url"),
            "slides_processed": len(slides)
        }

# --- Main Execution Block ---
if __name__ == "__main__":
    try:
        storage = CloudinaryStorage()
        
        # DEFINITIONS
        # Replace "test_presentation.pptx" with the actual name of your local file.
        LOCAL_PPTX_PATH = "C:/Users/Desktop/Ed-Tech/backend/teacher/media_toolkit/test_presentation.pptx" 
        
        # The ID it will have in Cloudinary (usually kept the same as filename for clarity)
        CLOUDINARY_PUBLIC_ID = "test_presentation.pptx" 

        # 1. Check if Local File Exists (Strict Check)
        if not os.path.exists(LOCAL_PPTX_PATH):
            logging.error(f"File not found: '{LOCAL_PPTX_PATH}'")
            logging.error("Please ensure you have placed your .pptx file in the correct directory.")
            exit(1) # Exit script if file is missing

        # 2. Upload to Cloudinary
        logging.info(f"Uploading local file '{LOCAL_PPTX_PATH}' to Cloudinary as '{CLOUDINARY_PUBLIC_ID}'...")
        success, _ = storage.upload_file(LOCAL_PPTX_PATH, CLOUDINARY_PUBLIC_ID)
        
        if success:
            # 3. Run Conversion
            converter = PPTXToHeyGenVideo(storage_manager=storage)
            # Pass the Cloudinary Public ID (not the local path) to the converter
            result = converter.convert(CLOUDINARY_PUBLIC_ID, title="Auto Generated Video")
            print("\n--- Conversion Successful ---")
            print(json.dumps(result, indent=2))
        else:
            logging.error("Could not upload file to Cloudinary. Aborting.")

    except Exception as e:
        logging.error(f"An error occurred: {e}", exc_info=True)