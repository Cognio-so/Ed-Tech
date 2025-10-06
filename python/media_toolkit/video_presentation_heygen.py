import os
import time
import json
import logging
import mimetypes
from typing import List, Dict, Optional, Any

import requests
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- OpenAI SDK Setup ---
try:
    from openai import OpenAI
except ImportError:
    raise ImportError("OpenAI SDK not found. Install with: pip install openai")

load_dotenv()

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


class PPTXToHeyGenVideo:
    def __init__(
        self,
        heygen_api_key: Optional[str] = None,
        openai_api_key: Optional[str] = None,
        avatar_id: Optional[str] = None,
        voice_id: Optional[str] = None,
        model: str = "gpt-4o",
        width: int = 1920,
        height: int = 1080,
        background_color: str = "#FFFFFF",
        request_timeout_s: int = 70,
        poll_interval_s: int = 5,
        pptx_avatar_id: Optional[str] = None,
        pptx_voice_id: Optional[str] = None,
        use_slides_as_background: bool = True,
        language: str = "english",  # NEW: Language support
    ):
        self.heygen_api_key = heygen_api_key or os.getenv("HEYGEN_API_KEY")
        self.openai_api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
        self.avatar_id = pptx_avatar_id or os.getenv("HEYGEN_AVATAR_ID")
        self.voice_id = pptx_voice_id or os.getenv("HEYGEN_VOICE_ID")
        self.model = model
        self.width = width
        self.height = height
        self.background_color = background_color
        self.request_timeout_s = request_timeout_s
        self.poll_interval_s = poll_interval_s
        self.use_slides_as_background = use_slides_as_background
        self.language = language.lower()  # NEW: Store language preference
        self.slide_asset_ids: List[str] = []

        if not self.heygen_api_key:
            raise ValueError("Missing HEYGEN_API_KEY")
        if not self.openai_api_key:
            raise ValueError("Missing OPENAI_API_KEY")
        if not self.avatar_id:
            raise ValueError("Missing HEYGEN_AVATAR_ID")
        if not self.voice_id:
            raise ValueError("Missing HEYGEN_VOICE_ID")

        self._client = OpenAI(api_key=self.openai_api_key)
        self._heygen_base_v2 = "https://api.heygen.com/v2"
        self._headers = {
            "Accept": "application/json",
            "Content-Type": "application/json",
            "X-Api-Key": self.heygen_api_key,
        }

    def _post_with_retry(self, url: str, payload: Dict) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.post(url, headers=self._headers, json=payload, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"POST attempt {attempt + 1} failed: {e}")
                if attempt == 2:
                    raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    def _get_with_retry(self, url: str) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.get(url, headers=self._headers, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"GET attempt {attempt + 1} failed: {e}")
                if attempt == 2:
                    raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    # NEW: Fetch actual voices from HeyGen API using correct endpoint
    def get_available_voices_from_api(self) -> List[Dict]:
        """Fetch actual voices from HeyGen API using correct endpoint"""
        try:
            # Try the correct endpoint
            url = "https://api.heygen.com/v1/brand_voice/list"
            response = self._get_with_retry(url)
            data = response.json()
            
            voices = data.get("data", [])
            logging.info(f"Fetched {len(voices)} voices from HeyGen API")
            
            # Log all voices to see what's available
            for voice in voices[:10]:  # Log first 10 voices
                voice_id = voice.get('voice_id') or voice.get('id')
                voice_name = voice.get('name', 'Unknown')
                language = voice.get('language', 'Unknown')
                logging.info(f"Voice: {voice_name} - ID: {voice_id} - Language: {language}")
            
            return voices
                
        except Exception as e:
            logging.error(f"Error fetching voices from API: {e}")
            return []

    # NEW: Find Arabic-compatible voices
    def find_arabic_voices(self) -> List[str]:
        """Find voices that actually support Arabic"""
        try:
            voices = self.get_available_voices_from_api()
            
            arabic_voices = []
            for voice in voices:
                voice_id = voice.get('voice_id') or voice.get('id')
                voice_name = voice.get('name', 'Unknown')
                language = voice.get('language', '')
                
                # Check if voice supports Arabic
                if (language.lower() == 'arabic' or 
                    'arabic' in language.lower() or
                    'ar' in language.lower() or
                    'arabic' in voice_name.lower()):
                    
                    arabic_voices.append(voice_id)
                    logging.info(f"Found Arabic voice: {voice_name} - ID: {voice_id}")
            
            if not arabic_voices:
                logging.warning("No Arabic voices found in API response")
                # Try some alternative voice IDs that might work
                alternative_voices = [
                    "bb2850ee8c76464d8e3d43f51b963fd1",  # Christine
                    "1776ddbd05374fa480e92f0297bbc67e",  # Melissa
                    "080f8e5cb3ae424989242b0efe5205e6",  # Ceecee
                    "baae7852b7824c8aaec62fc1c4e3064b",  # Rex
                ]
                logging.info(f"Using alternative voices: {alternative_voices}")
                return alternative_voices
            
            return arabic_voices
            
        except Exception as e:
            logging.error(f"Error finding Arabic voices: {e}")
            return ["bb2850ee8c76464d8e3d43f51b963fd1"]  # Fallback

    # NEW: Use actual Arabic voice IDs from voice_details.txt
    def get_voice_for_language(self, language: str) -> str:
        """Get the best voice for the specified language"""
        if language.lower() == 'arabic':
            # These are ACTUAL Arabic voice IDs from HeyGen
            arabic_voices = [
                "042173e02d18478384c64fdfe37ddd67",  # GHIZLANE (Female)
                "04fa555734714c3a90ac08a1ed64021c",  # Moncellence (Male)
                "0eb85e6e8710473b82f7e88609ba3053",  # Hushed Hiba - Excited (Female)
                "61cfb9ee298d419fa76d7f913f817447",  # Hakeem Hassan (Male)
                "61cfb9ee298d419fa76d7f913f817447",  # Sana (Female)
                "7042665eceec4300afd14e4f3ecf9157",  # Sana (Female)
                "e406a437e338443e9412162a0fff5289",  # Hushed Hiba (Female)
                "d12916aac1c44e6e8025ad820f1e9d4a",  # Hushed Hiba - Friendly (Female)
            ]
            
            # Use the first Arabic voice
            selected_voice = arabic_voices[0]
            logging.info(f"Selected Arabic voice: {selected_voice} (GHIZLANE)")
            return selected_voice
        else:
            return self.voice_id

    # NEW: Validate voice and language compatibility
    def validate_voice_language_compatibility(self, voice_id: str, language: str) -> bool:
        """Check if the selected voice supports the specified language"""
        try:
            # Since all our voices are multilingual, they support any language
            voices = self.get_available_voices_from_api(language)
            for voice in voices:
                if voice.get("voice_id") == voice_id:
                    return True
            return False
        except Exception as e:
            logging.error(f"Error validating voice compatibility: {e}")
            return False

    def _create_heygen_folder(self, folder_name: str) -> str:
        logging.info(f"Creating HeyGen asset folder: {folder_name}")
        payload = {"name": folder_name}
        url = "https://api.heygen.com/v1/folders/create"
        response = self._post_with_retry(url, payload)
        
        folder_id = response.json().get("data", {}).get("id")

        if not folder_id:
            raise RuntimeError(f"Failed to create HeyGen folder: {response.text}")
        logging.info(f"Successfully created folder with ID: {folder_id}")
        return folder_id

    def _upload_asset_to_heygen(self, file_path: str, folder_id: str) -> str:
        file_name = os.path.basename(file_path)
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            mime_type = "image/png"

        logging.info(f"Uploading {file_name} to HeyGen (target folder: {folder_id})...")

        url = "https://upload.heygen.com/v1/asset"
        
        upload_headers = {
            "X-Api-Key": self.heygen_api_key,
            "Content-Type": mime_type,
        }

        with open(file_path, 'rb') as f:
            for attempt in range(3):
                try:
                    f.seek(0)
                    resp = requests.post(url, headers=upload_headers, data=f, timeout=self.request_timeout_s)
                    resp.raise_for_status()
                    
                    response_data = resp.json()
                    asset_id = response_data.get("data", {}).get("id")

                    if not asset_id:
                        raise RuntimeError(f"Failed to get asset_id from HeyGen response: {response_data}")

                    logging.info(f"Successfully uploaded {file_name}. Asset ID: {asset_id}")
                    return asset_id

                except requests.RequestException as e:
                    logging.warning(f"Upload attempt {attempt + 1} for {file_name} failed: {e}")
                    if attempt == 2:
                        raise
                    time.sleep(2 ** attempt)
        
        raise RuntimeError(f"File upload failed for {file_name} after multiple retries.")

    def _extract_slides_as_images_and_upload(self, pptx_path: str, prs: Presentation, folder_id: str):
        import tempfile
        import sys
        
        temp_dir = tempfile.mkdtemp()
        logging.info(f"Created temporary directory for slide images: {temp_dir}")
        
        slide_image_paths = []
        try:
            if sys.platform == "win32":
                import win32com.client
                logging.info("Using PowerPoint COM automation to export slides...")
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True, WithWindow=False)
                
                num_slides_to_process = len(prs.slides)
                for i in range(1, num_slides_to_process + 1):
                    path = os.path.join(temp_dir, f"slide_{i}.png")
                    slide = presentation.Slides(i)
                    # FIXED: Use correct Export method parameters
                    slide.Export(path, "PNG")
                    slide_image_paths.append(path)
                    logging.info(f"Exported slide {i} to {path}")
                
                presentation.Close()
                powerpoint.Quit()
                logging.info("PowerPoint COM automation completed successfully")
                
            else:
                # For non-Windows systems, try using LibreOffice
                logging.info("Windows not detected, trying LibreOffice conversion...")
                slide_image_paths = self._convert_with_libreoffice(pptx_path, temp_dir)
                
        except Exception as e:
            logging.error(f"Slide export failed: {e}")
            raise RuntimeError(f"Failed to export slides from PowerPoint file: {e}")

        if not slide_image_paths:
            raise RuntimeError("Failed to extract any slide images")

        self.slide_asset_ids = [self._upload_asset_to_heygen(p, folder_id) for p in slide_image_paths]

    def _convert_with_libreoffice(self, pptx_path: str, temp_dir: str) -> List[str]:
        """Convert PowerPoint to images using LibreOffice"""
        try:
            import subprocess
            import os
            
            # Convert PPTX to PDF first
            pdf_path = os.path.join(temp_dir, "presentation.pdf")
            cmd = [
                "libreoffice", "--headless", "--convert-to", "pdf", 
                "--outdir", temp_dir, pptx_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")
            
            # Convert PDF to images using ImageMagick
            cmd = [
                "convert", "-density", "300", "-quality", "100",
                pdf_path, os.path.join(temp_dir, "slide_%d.png")
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode != 0:
                raise Exception(f"ImageMagick conversion failed: {result.stderr}")
            
            # Find generated PNG files
            slide_files = []
            for i in range(1, 100):  # Assume max 100 slides
                slide_path = os.path.join(temp_dir, f"slide_{i}.png")
                if os.path.exists(slide_path):
                    slide_files.append(slide_path)
                else:
                    break
            
            return slide_files
            
        except Exception as e:
            logging.warning(f"LibreOffice conversion failed: {e}")
            return []

    def _build_video_inputs(self, slide_notes: List[str]) -> List[Dict]:
        """FIXED: Proper background positioning and avatar placement with larger size"""
        video_inputs: List[Dict] = []
        
        for idx, note in enumerate(slide_notes):
            scene = {
                "character": {
                    "type": "talking_photo",
                    "talking_photo_id": self.avatar_id,
                    "talking_photo_style": "circle",
                    # FIXED: Increased avatar size to be more visible
                    "scale": 0.25,  # Increased from 0.08 to 0.25 (3x larger)
                    "offset": {
                        # FIXED: Position in bottom-right corner with proper margins
                        "x": 0.35,  # Move more to the right
                        "y": 0.35   # Move higher up to avoid slide content
                    },
                },
                "voice": {
                    "type": "text", 
                    "input_text": note, 
                    "voice_id": self.voice_id,
                    # NEW: Try adding language parameter to voice
                    "language": "ar" if self.language.lower() == 'arabic' else "en"
                },
            }

            # FIXED: Proper background handling - ensure slides are visible
            if self.use_slides_as_background and idx < len(self.slide_asset_ids):
                asset_id = self.slide_asset_ids[idx]
                logging.info(f"Setting background for scene {idx + 1} with asset_id: {asset_id}")
                scene["background"] = {
                    "type": "image", 
                    "image_asset_id": asset_id
                    # No position object - let HeyGen handle the full background
                }
            else:
                scene["background"] = {
                    "type": "color", 
                    "value": self.background_color
                }
            
            video_inputs.append(scene)
            
        return video_inputs

    def create_multi_scene_video(self, slide_notes: List[str], title: str) -> Dict:
        video_inputs = self._build_video_inputs(slide_notes)
        payload = {
            "video_inputs": video_inputs,
            "dimension": {"width": self.width, "height": self.height},
            "title": title,
        }
        generate_url = f"{self._heygen_base_v2}/video/generate"
        resp = self._post_with_retry(generate_url, payload)
        data = resp.json()
        video_id = data.get("data", {}).get("video_id")
        logging.info(f"Generating video using video_id : {video_id}, with Avatar id {self.avatar_id} and voice id {self.voice_id}.")
        if not video_id:
            raise RuntimeError(f"Unexpected response from generate endpoint: {resp.text}")
        return {"video_id": video_id}

    def wait_for_video(self, video_id: str) -> Dict:        
        """Polls the HeyGen API for the video status indefinitely until it is either 'completed' or 'failed'."""
        logging.info(f"Waiting for video {video_id} to complete. with Avatar id {self.avatar_id} and voice id {self.voice_id}. Polling status...")

        while True:
            status_url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"
            try:
                resp = self._get_with_retry(status_url)
                data = resp.json().get("data", {})
                status = data.get("status")

                if status in ("completed", "success"):
                    logging.info("Video generation completed successfully.")
                    return {"status": "completed", "video_url": data.get("video_url")}
                
                elif status in ("failed", "error"):
                    logging.error(f"Video generation failed with data: {data}")
                    raise RuntimeError(f"Video generation failed: {data}")
                
                else:
                    logging.info(f"Video status is '{status}', continuing to wait...")

            except requests.RequestException as e:
                logging.warning(f"Status check failed due to a network error: {e}. Retrying...")
            
            time.sleep(self.poll_interval_s)
    
    def generate_speaker_notes(self, slide_texts: List[str]) -> List[str]:
        logging.info(f"Generating speaker notes for {len(slide_texts)} slides using model {self.model}...")
        
        speaker_notes = []
        for i, text in enumerate(slide_texts):
            try:
                # FIXED: Language-aware system prompt
                system_prompt = (
                    f"You are a professional virtual teacher creating engaging educational content in {self.language}. "
                    "Your task is to explain presentation slide content in a clear, engaging, and educational manner. "
                    "Create speaker notes that are conversational, informative, and easy to follow. "
                    "Keep the explanation concise but comprehensive, suitable for video narration. "
                    f"Use a professional yet approachable tone in {self.language} that keeps viewers engaged."
                )
                
                user_prompt = f"Here is the slide content:\n\n---\n\n{text}\n\n---\n\nPlease provide engaging speaker notes for video narration in {self.language}. Keep it under 200 words and make it conversational."
                
                response = self._client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    temperature=0.7,
                    max_tokens=300,
                )
                
                note = response.choices[0].message.content.strip()
                speaker_notes.append(note)
                logging.info(f"Generated note for slide {i+1}: '{note[:100]}...'")

            except Exception as e:
                logging.error(f"Failed to generate notes for slide {i+1}: {e}")
                speaker_notes.append(f"Let me explain this slide content: {text}")

        return speaker_notes

    def convert(self, pptx_path: str, title: Optional[str] = None, max_slides: Optional[int] = None) -> Dict[str, Any]:
        """
        Convert a PowerPoint presentation to a HeyGen video.
        """
        logging.info(f"Starting conversion of {pptx_path}")
        
        # FIXED: Auto-select voice based on language
        if self.language.lower() == 'arabic':
            self.voice_id = self.get_voice_for_language(self.language)
            logging.info(f"Auto-selected voice for Arabic: {self.voice_id}")
        
        logging.info(f"Using voice {self.voice_id} for language {self.language}")
        
        prs = Presentation(pptx_path)
        slides = list(prs.slides)
        
        if not slides:
            raise ValueError("No slides found in the PowerPoint file.")

        if max_slides and len(slides) > max_slides:
            logging.warning(f"Limiting to {max_slides} slides from {len(slides)}.")
            slides = slides[:max_slides]
        
        # Better text extraction
        slides_text = []
        for slide in slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            slides_text.append("\n".join(slide_text) if slide_text else f"Slide {len(slides_text) + 1}")

        if self.use_slides_as_background:
            folder_name = f"Slides_{os.path.basename(pptx_path)}_{int(time.time())}"
            folder_id = self._create_heygen_folder(folder_name)
            self._extract_slides_as_images_and_upload(pptx_path, prs, folder_id)
            if len(self.slide_asset_ids) < len(slides_text):
                logging.warning("Number of uploaded slide assets is less than the number of slides.")
                slides_text = slides_text[:len(self.slide_asset_ids)]

        slide_notes = self.generate_speaker_notes(slides_text)
        
        final_title = title or os.path.basename(pptx_path)
        gen_info = self.create_multi_scene_video(slide_notes, title=final_title)
        final_status = self.wait_for_video(gen_info["video_id"])

        return {
            "video_id": gen_info["video_id"],
            "video_url": final_status.get("video_url"),
            "slides_count": len(slide_notes),
            "language": self.language,
            "voice_id": self.voice_id
        }

def main():
    pptx_file_path = input("enter ppt file path here :")
    pptx_avatar_id = input("enter ppt avatar id here :")
    pptx_voice_id = input("enter ppt voice id here :")
    language = input("enter language (english/arabic): ").strip().lower() or "english"

    if not os.path.exists(pptx_file_path):
        logging.error(f"File not found: {pptx_file_path}")
        logging.error("Please update the 'pptx_file_path' variable in the main() function with the correct path.")
        return

    try:
        converter = PPTXToHeyGenVideo(
            pptx_avatar_id=pptx_avatar_id,
            pptx_voice_id=pptx_voice_id,
            language=language
        )
        
        result = converter.convert(pptx_file_path)
        
        print("\n--- Conversion Successful ---")
        print(json.dumps(result, indent=2))
    except Exception as e:
        logging.error(f"An error occurred during conversion: {e}", exc_info=True)

if __name__ == "__main__":
    main()